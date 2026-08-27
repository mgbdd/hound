import subprocess, zipfile
from pathlib import Path
from pptx import Presentation
from pptx.exc import PackageNotFoundError
from tempfile import TemporaryDirectory
from processor_server.analyzers.image_file_analyzer import image_file_analyzer
from processor_server.analyzers.audio_file_analyzer import audio_file_analyzer
from processor_server.finder_soffice import find_soffice


def pptx_file_analyzer(file_path, img_limit=None, audio_limit=None):
    src = Path(file_path)

    with TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        try:
            if src.suffix.lower() == ".ppt":
                soffice_bin = find_soffice()
                subprocess.check_call([
                    soffice_bin, "--headless",
                    "--convert-to", "pptx",
                    "--outdir", tmpdir,
                    src
                ])
                src = tmpdir / f"{src.stem}.pptx"
                if not src.exists():
                    return ""
        except (FileNotFoundError, subprocess.CalledProcessError) as e:
            return ""

        try:
            prs = Presentation(str(src))
        except (PackageNotFoundError, zipfile.BadZipFile, KeyError) as e:

            return ""

        img_count = 0
        audio_count = 0
        text_out = []

        try:
            with zipfile.ZipFile(src, 'r') as z:
                for name in z.namelist():
                    if audio_limit is not None and audio_count >= audio_limit:
                        break
                    if name.startswith("ppt/media/") and name.lower().endswith((".mp3", ".wav")):
                        z.extract(name, path=tmpdir)
                        media_path = tmpdir / name
                        try:
                            desc = audio_file_analyzer(media_path) or "нет описания"
                            text_out.append(f"[audio {audio_count + 1}]({media_path.name}): {desc}")
                            audio_count += 1
                        except Exception as e:
                            continue
                            # logging.exception("audio_file_analyzer упал на %s: %s", media_path, e)
        except Exception as e:
            pass
            # logging.exception("Ошибка распаковки аудио из PPTX: %s", e)
        for slide_idx, slide in enumerate(prs.slides, start=1):

            for shape in slide.shapes:
                if getattr(shape, "text", "").strip():
                    text_out.append(shape.text.strip())

            for shape in slide.shapes:
                if shape.shape_type == 13 and (img_limit is None or img_count < img_limit):
                    try:
                        img = shape.image
                        ext = img.ext
                        img_name = f"slide{slide_idx}_img{img_count + 1}.{ext}"
                        img_path = tmpdir / img_name
                        img_path.write_bytes(img.blob)

                        alt = image_file_analyzer(img_path) or "нет описания"
                        text_out.append(f"![image {img_count + 1}]({img_name}): {alt}")
                        img_count += 1
                    except Exception as e:
                        continue
                        # logging.exception("image_file_analyzer упал на слайде %s: %s", slide_idx, e)

        return "".join(text_out)
